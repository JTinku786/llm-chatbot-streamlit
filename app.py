"""
Modern ChatGPT/Perplexity-style LLM Application with Multi-Chat & File Upload
Built with OpenAI, Pinecone, and Advanced Document Processing
"""

import os

import streamlit as st
from langsmith import traceable
from langsmith.wrappers import wrap_openai
from openai  import OpenAI
from PIL import Image
import base64
from io import BytesIO
from datetime import datetime
import uuid
import json
import re
import requests
from pypdf import PdfReader
from pptx import Presentation
import docx

# Page configuration
st.set_page_config(
    page_title="AI Chat Assistant",
    page_icon="🤖",
    layout="wide",
    initial_sidebar_state="expanded"
)

# Custom CSS for ChatGPT-style UI
st.markdown("""
<style>
    #MainMenu {visibility: hidden;}
    footer {visibility: hidden;}
    header {visibility: hidden;}
    
    .main {
        background-color: #f7f7f8;
    }
    
    .stChatMessage {
        padding: 1rem;
        margin: 0.5rem 0;
        border-radius: 0.5rem;
    }
    
    .stChatInputContainer {
        border-top: 1px solid #ececec;
        padding: 1rem 0;
    }
    
    [data-testid="stSidebar"] {
        background-color: #202123;
    }
    
    [data-testid="stSidebar"] .stMarkdown {
        color: #ececec;
    }
    
    [data-testid="stSidebar"] .stButton button {
        background-color: transparent;
        color: #ececec;
        border: 1px solid #565869;
        width: 100%;
        text-align: left;
        padding: 0.5rem 1rem;
        border-radius: 0.5rem;
        margin: 0.25rem 0;
    }
    
    [data-testid="stSidebar"] .stButton button:hover {
        background-color: #2a2b32;
    }
    
    h1 {
        font-weight: 600;
        font-size: 2rem;
        color: #202123;
        text-align: center;
        padding: 2rem 0 1rem 0;
    }
    
    .block-container {
        padding-top: 2rem;
        padding-bottom: 0rem;
        max-width: 48rem;
    }
    
    .file-upload-area {
        border: 2px dashed #565869;
        border-radius: 0.5rem;
        padding: 1rem;
        margin: 1rem 0;
        background-color: #f7f7f8;
        text-align: center;
    }
    
    .uploaded-file {
        background-color: #e3f2fd;
        padding: 0.5rem 1rem;
        border-radius: 0.5rem;
        margin: 0.5rem 0;
        display: inline-block;
    }
</style>
""", unsafe_allow_html=True)

# Load configuration
@st.cache_resource
def load_config():
    try:
        return {
            "openai_api_key": st.secrets["OPENAI_API_KEY"],
            "pinecone_api_key": st.secrets.get("PINECONE_API_KEY", os.getenv("PINECONE_API_KEY", "")),
            "pinecone_environment": st.secrets.get("PINECONE_ENVIRONMENT", "gcp-starter"),
            "pinecone_index_name": st.secrets.get("PINECONE_INDEX_NAME", os.getenv("PINECONE_INDEX_NAME", "chatbot-memory")),
            "langsmith_api_key": st.secrets.get("LANGSMITH_API_KEY", ""),
            "langsmith_project": st.secrets.get("LANGSMITH_PROJECT", "llm-chatbot-streamlit"),
            "langsmith_endpoint": st.secrets.get("LANGSMITH_ENDPOINT", "https://api.smith.langchain.com"),
            "openweathermap_api_key": st.secrets.get("OPENWEATHERMAP_API_KEY", ""),
            "serpapi_api_key": st.secrets.get("SERPAPI_API_KEY", ""),
            "tavily_api_key": st.secrets.get("TAVILY_API_KEY", ""),
            "ict_knowledge_index_name": st.secrets.get("ICT_KNOWLEDGE_INDEX", st.secrets.get("PINECONE_INDEX_NAME", os.getenv("PINECONE_INDEX_NAME", "ict-knowledge-index"))),
            "chat_memory_index_name": st.secrets.get("CHAT_MEMORY_INDEX", os.getenv("CHAT_MEMORY_INDEX", "chat-memory-index")),
            "ict_domain": st.secrets.get("ICT_DOMAIN", "ict_trading"),
            "retrieval_score_threshold": float(st.secrets.get("RETRIEVAL_SCORE_THRESHOLD", 0.3)),
        }
    except Exception as e:
        st.error(f"Error loading configuration: {e}")
        return None

config = load_config()
if not config:
    st.stop()

if config["langsmith_api_key"]:
    os.environ["LANGSMITH_API_KEY"] = config["langsmith_api_key"]
    os.environ["LANGSMITH_TRACING"] = "true"
    os.environ["LANGSMITH_PROJECT"] = config["langsmith_project"]
    os.environ["LANGSMITH_ENDPOINT"] = config["langsmith_endpoint"]

client = wrap_openai(OpenAI(api_key=config["openai_api_key"]))

@traceable(name="chat_completion_stream", run_type="llm")
def stream_chat_completion(selected_model, messages, temperature):
    return client.chat.completions.create(
        model=selected_model,
        messages=messages,
        temperature=temperature,
        stream=True
    )









@traceable(name="rewrite_ict_query", run_type="tool")
def rewrite_ict_query(query, model="gpt-4.1-mini"):
    """Rewrite ambiguous ICT questions to retrieval-friendly form."""
    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "Rewrite the user query for ICT (Inner Circle Trader) document retrieval. Keep intent unchanged, add disambiguating ICT terms if vague, and return only one concise rewritten query."},
                {"role": "user", "content": query},
            ],
            temperature=0.0,
            max_tokens=80,
        )
        rewritten = (resp.choices[0].message.content or "").strip()
        return rewritten or query
    except Exception:
        return query


@traceable(name="summarize_last_5_turns", run_type="tool")
def summarize_last_5_turns(messages, model="gpt-4.1-mini"):
    """Controlled memory injection: summarize only last 5 turns."""
    turns = []
    for msg in messages:
        role = msg.get("role")
        if role not in {"user", "assistant"}:
            continue
        content = msg.get("content", "")
        if isinstance(content, list):
            text_parts = [p.get("text", "") for p in content if isinstance(p, dict) and p.get("type") == "text"]
            content = "\n".join(text_parts)
        turns.append(f"{role}: {str(content)[:1000]}")

    if not turns:
        return ""

    window = turns[-10:]
    try:
        resp = client.chat.completions.create(
            model=model,
            messages=[
                {"role": "system", "content": "Summarize the conversation memory in <=8 bullet points. Preserve facts, open tasks, constraints, and user preferences. Exclude unrelated topics."},
                {"role": "user", "content": "\n".join(window)},
            ],
            temperature=0.0,
            max_tokens=220,
        )
        return (resp.choices[0].message.content or "").strip()
    except Exception:
        return ""


def resolve_pinecone_index(index_name):
    """Resolve Pinecone index host + dimension using REST API."""
    if not config["pinecone_api_key"]:
        return None, None, "PINECONE_API_KEY is not configured.", None

    headers = {"Api-Key": config["pinecone_api_key"], "Accept": "application/json"}
    try:
        response = requests.get(
            f"https://api.pinecone.io/indexes/{index_name}",
            headers=headers,
            timeout=15,
        )
        if response.status_code == 404:
            return None, None, f"Pinecone index '{index_name}' does not exist.", None
        response.raise_for_status()
        data = response.json()
        host = data.get("host") or data.get("status", {}).get("host")
        dimension = data.get("dimension")
        if dimension is None:
            dimension = data.get("spec", {}).get("dimension")
        if not host:
            return None, None, "Pinecone index host could not be resolved.", None
        return host, headers, "", dimension
    except requests.RequestException as exc:
        return None, None, f"Pinecone index lookup failed: {exc}", None


def extract_search_query(prompt):
    """Extract web-search query from user prompt."""
    prompt_clean = prompt.strip()
    prompt_lower = prompt_clean.lower()

    if prompt_lower.startswith('/search '):
        return prompt_clean[len('/search '):].strip()
    if prompt_lower.startswith('search '):
        return prompt_clean[len('search '):].strip()
    if prompt_lower.startswith('google '):
        return prompt_clean[len('google '):].strip()
    if prompt_lower.startswith('look up '):
        return prompt_clean[len('look up '):].strip()

    return ''


@traceable(name="search_serpapi", run_type="tool")
def search_with_serpapi(query, max_results=5):
    """Fetch search results from SerpAPI."""
    if not config["serpapi_api_key"]:
        return [], "SERPAPI_API_KEY is not configured."

    try:
        response = requests.get(
            "https://serpapi.com/search.json",
            params={
                "q": query,
                "api_key": config["serpapi_api_key"],
                "num": max_results,
            },
            timeout=12,
        )
        response.raise_for_status()
        data = response.json()
        items = data.get("organic_results", [])[:max_results]
        results = [
            {
                "title": item.get("title", "Untitled"),
                "snippet": item.get("snippet", ""),
                "url": item.get("link", ""),
            }
            for item in items
        ]
        return results, ""
    except requests.RequestException as exc:
        return [], f"SerpAPI search failed: {exc}"


@traceable(name="search_tavily", run_type="tool")
def search_with_tavily(query, max_results=5):
    """Fetch search results from Tavily."""
    if not config["tavily_api_key"]:
        return [], "TAVILY_API_KEY is not configured."

    try:
        response = requests.post(
            "https://api.tavily.com/search",
            json={
                "api_key": config["tavily_api_key"],
                "query": query,
                "max_results": max_results,
                "include_answer": True,
            },
            timeout=12,
        )
        response.raise_for_status()
        data = response.json()
        items = data.get("results", [])[:max_results]
        results = [
            {
                "title": item.get("title", "Untitled"),
                "snippet": item.get("content", ""),
                "url": item.get("url", ""),
            }
            for item in items
        ]
        if data.get("answer"):
            results.insert(0, {"title": "Tavily Answer", "snippet": data["answer"], "url": ""})
        return results[:max_results], ""
    except requests.RequestException as exc:
        return [], f"Tavily search failed: {exc}"


@traceable(name="web_search_context", run_type="tool")
def load_web_search_context(query, provider):
    """Load web search context from selected provider."""
    provider_order = {
        "Auto": ["Tavily", "SerpAPI"],
        "Tavily": ["Tavily"],
        "SerpAPI": ["SerpAPI"],
    }

    provider_map = {
        "Tavily": search_with_tavily,
        "SerpAPI": search_with_serpapi,
    }

    errors = []
    for provider_name in provider_order.get(provider, [provider]):
        results, error = provider_map[provider_name](query)
        if results:
            lines = [
                f"[{idx}] {item['title']}\n{item['snippet']}\n{item['url']}"
                for idx, item in enumerate(results, start=1)
            ]
            return "\n\n".join(lines), provider_name, ""
        if error:
            errors.append(f"{provider_name}: {error}")

    return "", "", " | ".join(errors) if errors else "No search results found."


def extract_weather_cities(prompt):
    """Extract city names from weather-related prompts."""
    prompt_clean = prompt.strip()
    prompt_lower = prompt_clean.lower()

    if prompt_lower.startswith('/weather '):
        city_fragment = prompt_clean[len('/weather '):]
    elif 'weather in ' in prompt_lower:
        start = prompt_lower.find('weather in ') + len('weather in ')
        city_fragment = prompt_clean[start:]
    elif 'temperature in ' in prompt_lower:
        start = prompt_lower.find('temperature in ') + len('temperature in ')
        city_fragment = prompt_clean[start:]
    elif 'forecast in ' in prompt_lower:
        start = prompt_lower.find('forecast in ') + len('forecast in ')
        city_fragment = prompt_clean[start:]
    elif prompt_lower.startswith('weather '):
        city_fragment = prompt_clean[len('weather '):]
    else:
        return []

    city_fragment = re.split(r'[?.!\n]', city_fragment, maxsplit=1)[0]
    city_fragment = city_fragment.strip()
    if not city_fragment:
        return []

    city_fragment = re.sub(r'\b(today|now|currently|please)\b', '', city_fragment, flags=re.IGNORECASE).strip()
    parts = re.split(r',|\band\b|&', city_fragment, flags=re.IGNORECASE)
    return [part.strip(' .') for part in parts if part.strip(' .')]


@traceable(name="weather_context", run_type="tool")
def load_weather_context(cities):
    """Load weather context for the given cities using OpenWeatherMap."""
    if not config["openweathermap_api_key"]:
        return "", "OPENWEATHERMAP_API_KEY is not configured."

    clean_cities = [city.strip() for city in cities if city and city.strip()]
    if not clean_cities:
        return "", "No valid city names provided."

    try:
        weather_reports = []
        for city in clean_cities:
            response = requests.get(
                "https://api.openweathermap.org/data/2.5/weather",
                params={
                    "q": city,
                    "appid": config["openweathermap_api_key"],
                    "units": "metric",
                },
                timeout=10,
            )
            response.raise_for_status()
            data = response.json()

            weather_reports.append(
                f"{data['name']}: {data['weather'][0]['description']}, "
                f"temperature {data['main']['temp']}°C, feels like {data['main']['feels_like']}°C, "
                f"humidity {data['main']['humidity']}%, wind {data['wind']['speed']} m/s"
            )

        return "\n".join(weather_reports), ""
    except requests.RequestException as exc:
        return "", f"Weather lookup failed: {exc}"


@traceable(name="tool_router", run_type="chain")
def route_tools(prompt, provider):
    """Decide whether tools are needed and fetch context before LLM call."""
    routing = {
        "weather_context": "",
        "weather_error": "",
        "weather_cities": [],
        "search_context": "",
        "search_error": "",
        "search_query": "",
        "provider_used": "",
    }

    city_list = extract_weather_cities(prompt)
    if city_list:
        routing["weather_cities"] = city_list
        routing["weather_context"], routing["weather_error"] = load_weather_context(city_list)

    explicit_search_query = extract_search_query(prompt)
    time_sensitive_pattern = r"\b(as of|latest|current|today|breaking|recent|news|update|updates|happening now)\b"
    needs_fresh_info = bool(re.search(time_sensitive_pattern, prompt, flags=re.IGNORECASE))

    if explicit_search_query:
        routing["search_query"] = explicit_search_query
    elif needs_fresh_info and not city_list:
        routing["search_query"] = prompt.strip()

    if routing["search_query"]:
        (
            routing["search_context"],
            routing["provider_used"],
            routing["search_error"],
        ) = load_web_search_context(routing["search_query"], provider)

    return routing


@traceable(name="ict_retrieve_chunks", run_type="tool")
def retrieve_ict_chunks(query, top_k=8, domain_filter="ict_trading", source_filter="", trust_tier="approved", score_threshold=0.3):
    """Dense retrieval from ICT knowledge index with production filters + confidence gating."""
    index_host, headers, init_error, index_dimension = resolve_pinecone_index(config["ict_knowledge_index_name"])
    if not index_host:
        return [], {"error": init_error}

    try:
        emb_params = {"model": "text-embedding-3-small", "input": [query], "encoding_format": "float"}
        if index_dimension and int(index_dimension) <= 1536:
            emb_params["dimensions"] = int(index_dimension)
        q_emb = client.embeddings.create(**emb_params).data[0].embedding

        and_filters = [{"source": {"$ne": "streamlit_chat"}}]
        if domain_filter:
            and_filters.append({"domain": {"$eq": domain_filter}})
        if source_filter:
            and_filters.append({"source": {"$eq": source_filter}})
        if trust_tier and trust_tier.lower() != "all":
            and_filters.append({"trust_tier": {"$eq": trust_tier}})

        query_body = {
            "vector": q_emb,
            "topK": top_k,
            "includeMetadata": True,
            "filter": {"$and": and_filters},
        }

        query_resp = requests.post(
            f"https://{index_host}/query",
            headers={**headers, "Content-Type": "application/json"},
            json=query_body,
            timeout=20,
        )
        if query_resp.status_code >= 400:
            return [], {"error": f"Pinecone query failed HTTP {query_resp.status_code}: {query_resp.text[:800]}"}

        matches = query_resp.json().get("matches", [])
        all_docs = []
        for m in matches:
            meta = m.get("metadata", {}) or {}
            all_docs.append({
                "source_id": meta.get("source_id", ""),
                "doc_id": meta.get("doc_id", ""),
                "source": meta.get("source", "unknown"),
                "domain": meta.get("domain", ""),
                "trust_tier": meta.get("trust_tier", ""),
                "chunk_id": meta.get("chunk_id", ""),
                "text": meta.get("text", ""),
                "score": float(m.get("score", 0)),
            })

        accepted = [d for d in all_docs if d["score"] >= score_threshold]
        rejected = len(all_docs) - len(accepted)
        info = {
            "error": "",
            "filters": {
                "domain": domain_filter or "any",
                "source": source_filter or "any",
                "trust_tier": trust_tier or "any",
                "exclude_source": "streamlit_chat",
            },
            "score_threshold": score_threshold,
            "retrieved_total": len(all_docs),
            "accepted_total": len(accepted),
            "rejected_low_quality": rejected,
            "index_name": config["ict_knowledge_index_name"],
        }
        if not accepted:
            info["error"] = "No high-confidence ICT chunks after filtering/threshold."
        return accepted, info
    except requests.RequestException as exc:
        return [], {"error": f"ICT retrieval failed: {exc}"}


def build_ict_rag_context(query, top_k=8, domain_filter="ict_trading", source_filter="", trust_tier="approved", score_threshold=0.3, rewrite_model="gpt-4.1-mini"):
    """Build filtered RAG context with retrieval diagnostics."""
    rewritten_query = rewrite_ict_query(query, model=rewrite_model)
    docs, info = retrieve_ict_chunks(
        rewritten_query,
        top_k=top_k,
        domain_filter=domain_filter,
        source_filter=source_filter,
        trust_tier=trust_tier,
        score_threshold=score_threshold,
    )
    err = info.get("error", "")
    if err:
        return "", [], info, rewritten_query

    parts = []
    for i, d in enumerate(docs, start=1):
        parts.append(
            f"### Chunk {i} (doc_id: {d['doc_id']}, chunk_id: {d['chunk_id']}, source: {d['source']}, domain: {d['domain']}, trust_tier: {d['trust_tier']}, score: {d['score']:.3f})\n"
            f"{d['text'][:1200]}"
        )

    return "\n\n".join(parts), docs, info, rewritten_query


def get_pinecone_index(index_name=None):
    """Backward-compatible wrapper for chat memory index resolution."""
    return resolve_pinecone_index(index_name or config["chat_memory_index_name"])


@traceable(name="store_conversation_pinecone", run_type="tool")
def store_conversation_in_pinecone(chat_id, user_message, assistant_message, summary=""):
    """Store chat memory only in chat-memory index (never in ICT knowledge index)."""
    index_host, headers, init_error, index_dimension = get_pinecone_index(config["chat_memory_index_name"])
    if not index_host:
        return {
            "success": False,
            "reason": init_error,
            "vector_id": "",
            "error_type": "init_error",
            "index_dimension": index_dimension,
            "index_name": config["chat_memory_index_name"],
        }

    text_payload = f"User: {user_message}\nAssistant: {assistant_message}\nSummary: {summary}"

    try:
        embedding_params = {"model": "text-embedding-3-small", "input": text_payload}
        if index_dimension:
            if int(index_dimension) > 1536:
                return {
                    "success": False,
                    "reason": f"Index dimension {index_dimension} is larger than text-embedding-3-small max 1536.",
                    "vector_id": "",
                    "error_type": "dimension_mismatch",
                    "index_dimension": index_dimension,
                    "index_name": config["chat_memory_index_name"],
                }
            embedding_params["dimensions"] = int(index_dimension)

        embedding = client.embeddings.create(**embedding_params).data[0].embedding

        turn_id = str(uuid.uuid4())
        vector_id = f"chat-{chat_id}-{turn_id}"
        metadata = {
            "source_id": f"chat_{chat_id}",
            "doc_id": f"chat_session_{chat_id}",
            "chunk_id": f"turn_{turn_id}",
            "chat_id": chat_id,
            "timestamp": datetime.utcnow().isoformat(),
            "user_message": user_message[:2000],
            "assistant_message": assistant_message[:2000],
            "conversation_summary": summary[:2000],
            "source": "streamlit_chat",
            "domain": "chat_memory",
            "trust_tier": "user_generated",
        }

        upsert_payload = {"vectors": [{"id": vector_id, "values": embedding, "metadata": metadata}]}
        upsert_response = requests.post(
            f"https://{index_host}/vectors/upsert",
            headers={**headers, "Content-Type": "application/json"},
            json=upsert_payload,
            timeout=15,
        )

        if upsert_response.status_code >= 400:
            error_body = upsert_response.text[:1000]
            return {
                "success": False,
                "reason": f"Pinecone REST upsert failed: HTTP {upsert_response.status_code}. Response: {error_body}",
                "vector_id": "",
                "error_type": "HTTPError",
                "index_dimension": index_dimension,
                "embedding_length": len(embedding),
                "index_name": config["chat_memory_index_name"],
            }

        return {
            "success": True,
            "reason": "",
            "vector_id": vector_id,
            "error_type": "",
            "index_dimension": index_dimension,
            "embedding_length": len(embedding),
            "index_name": config["chat_memory_index_name"],
        }
    except requests.RequestException as exc:
        return {
            "success": False,
            "reason": f"Pinecone REST upsert failed: {exc}",
            "vector_id": "",
            "error_type": type(exc).__name__,
            "index_dimension": index_dimension,
            "embedding_length": 0,
            "index_name": config["chat_memory_index_name"],
        }
    except Exception as exc:
        return {
            "success": False,
            "reason": f"Pinecone upsert failed: {exc}",
            "vector_id": "",
            "error_type": type(exc).__name__,
            "index_dimension": index_dimension,
            "index_name": config["chat_memory_index_name"],
        }


# Helper functions for file processing
def encode_image_to_base64(image):
    """Convert PIL Image to base64 string."""
    buffered = BytesIO()
    image.save(buffered, format="PNG")
    return base64.b64encode(buffered.getvalue()).decode()

def extract_text_from_pdf(pdf_file):
    """Extract text from PDF file."""
    try:
        pdf_reader = PdfReader(pdf_file)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text() + "\n"
        return text
    except Exception as e:
        return f"Error extracting PDF: {str(e)}"

def extract_text_from_pptx(pptx_file):
    """Extract text from PowerPoint file."""
    try:
        prs = Presentation(pptx_file)
        text = ""
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text += shape.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting PPTX: {str(e)}"

def extract_text_from_docx(docx_file):
    """Extract text from Word document."""
    try:
        doc = docx.Document(docx_file)
        text = ""
        for para in doc.paragraphs:
            text += para.text + "\n"
        return text
    except Exception as e:
        return f"Error extracting DOCX: {str(e)}"

def process_uploaded_file(uploaded_file):
    """Process uploaded file and return content."""
    file_type = uploaded_file.type
    file_name = uploaded_file.name
    
    if file_type.startswith('image/'):
        # Image processing
        image = Image.open(uploaded_file)
        return {
            "type": "image",
            "name": file_name,
            "content": image,
            "base64": encode_image_to_base64(image)
        }
    elif file_type == 'application/pdf':
        # PDF processing
        text = extract_text_from_pdf(uploaded_file)
        return {
            "type": "pdf",
            "name": file_name,
            "content": text
        }
    elif file_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
        # PowerPoint processing
        text = extract_text_from_pptx(uploaded_file)
        return {
            "type": "pptx",
            "name": file_name,
            "content": text
        }
    elif file_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
        # Word document processing
        text = extract_text_from_docx(uploaded_file)
        return {
            "type": "docx",
            "name": file_name,
            "content": text
        }
    else:
        return {
            "type": "unknown",
            "name": file_name,
            "content": "Unsupported file type"
        }

# Initialize session state
if "all_chats" not in st.session_state:
    st.session_state.all_chats = {}

if "current_chat_id" not in st.session_state:
    chat_id = datetime.now().strftime("%Y%m%d_%H%M%S")
    st.session_state.current_chat_id = chat_id
    st.session_state.all_chats[chat_id] = {
        "title": "New Chat",
        "messages": [],
        "created_at": datetime.now().isoformat()
    }

if "uploaded_files" not in st.session_state:
    st.session_state.uploaded_files = []

# Helper functions for chat management
def create_new_chat():
    chat_id = datetime.now().strftime("%Y%m%d_%H%M%S_%f")
    st.session_state.all_chats[chat_id] = {
        "title": "New Chat",
        "messages": [],
        "created_at": datetime.now().isoformat()
    }
    st.session_state.current_chat_id = chat_id
    st.session_state.uploaded_files = []

def delete_chat(chat_id):
    if len(st.session_state.all_chats) > 1:
        del st.session_state.all_chats[chat_id]
        st.session_state.current_chat_id = list(st.session_state.all_chats.keys())[-1]
    else:
        st.session_state.all_chats[chat_id]["messages"] = []
        st.session_state.all_chats[chat_id]["title"] = "New Chat"

def switch_chat(chat_id):
    st.session_state.current_chat_id = chat_id
    st.session_state.uploaded_files = []

def get_chat_title(messages):
    if messages:
        first_user_msg = next((msg["content"] for msg in messages if msg["role"] == "user"), None)
        if first_user_msg and isinstance(first_user_msg, str):
            return first_user_msg[:50] + "..." if len(first_user_msg) > 50 else first_user_msg
    return "New Chat"

current_chat = st.session_state.all_chats[st.session_state.current_chat_id]

# Sidebar
with st.sidebar:
    st.markdown("### 🤖 AI Chat Assistant")
    
    if st.button("➕ New Chat", use_container_width=True, key="new_chat_btn"):
        create_new_chat()
        st.rerun()
    
    st.markdown("---")
    st.markdown("### 💬 Chat History")
    
    sorted_chats = sorted(
        st.session_state.all_chats.items(),
        key=lambda x: x[1]["created_at"],
        reverse=True
    )
    
    for chat_id, chat_data in sorted_chats:
        if chat_data["messages"]:
            chat_data["title"] = get_chat_title(chat_data["messages"])
        
        col1, col2 = st.columns([5, 1])
        
        with col1:
            is_active = chat_id == st.session_state.current_chat_id
            button_type = "primary" if is_active else "secondary"
            
            if st.button(
                f"{'📍 ' if is_active else '💬 '}{chat_data['title']}",
                key=f"chat_{chat_id}",
                use_container_width=True,
                type=button_type
            ):
                switch_chat(chat_id)
                st.rerun()
        
        with col2:
            if st.button("🗑️", key=f"del_{chat_id}", help="Delete chat"):
                delete_chat(chat_id)
                st.rerun()
    
    st.markdown("---")
    st.markdown("### ⚙️ Settings")
    
    selected_model = st.selectbox(
        "Model",
        ["gpt-4o", "gpt-4o-mini", "gpt-4-turbo", "gpt-3.5-turbo"],
        index=0  # gpt-4o for vision support
    )
    
    temperature = st.slider("Temperature", 0.0, 2.0, 0.7, 0.1)
    use_rag = st.toggle("Use Memory (RAG)", value=True)
    web_search_provider = st.selectbox(
        "Web Search Provider",
        ["Auto", "Tavily", "SerpAPI"],
        index=0,
        help="Used when you ask with /search, search, google, or look up."
    )

    st.markdown("### 📚 ICT Concept RAG")
    use_ict_rag = st.checkbox("Include ICT Concept Technology (RAG)", value=False)
    rag_domain_filter = st.text_input("RAG Domain Filter", value=config["ict_domain"], help="Example: ict_trading")
    rag_source_filter = st.text_input("RAG Source Filter", value="", help="Optional source filter, e.g., ict_manual")
    rag_trust_tier = st.selectbox("RAG Trust Tier", ["approved", "all"], index=0)

    st.markdown("---")
    st.markdown("### 📊 Stats")
    st.metric("Total Chats", len(st.session_state.all_chats))
    st.metric("Messages", len(current_chat["messages"]))
    st.metric("Files Attached", len(st.session_state.uploaded_files))

# Main chat interface
st.title("💬 AI Chat Assistant")

# File upload section
with st.expander("📎 Attach Files (Images, PDFs, PPT, Docs)", expanded=bool(st.session_state.uploaded_files)):
    uploaded_files = st.file_uploader(
        "Upload files to analyze",
        type=["png", "jpg", "jpeg", "pdf", "pptx", "docx"],
        accept_multiple_files=True,
        key="file_uploader"
    )
    
    if uploaded_files:
        st.session_state.uploaded_files = []
        for uploaded_file in uploaded_files:
            processed = process_uploaded_file(uploaded_file)
            st.session_state.uploaded_files.append(processed)
            
            # Display uploaded file
            if processed["type"] == "image":
                col1, col2 = st.columns([1, 4])
                with col1:
                    st.image(processed["content"], width=100)
                with col2:
                    st.success(f"✅ {processed['name']}")
            else:
                st.success(f"✅ {processed['name']} ({processed['type'].upper()}) - {len(processed['content'])} characters extracted")

# Display chat messages
for message in current_chat["messages"]:
    with st.chat_message(message["role"]):
        if isinstance(message["content"], list):
            # Handle multimodal content
            for content_part in message["content"]:
                if isinstance(content_part, dict):
                    if content_part.get("type") == "image_url":
                        st.image(content_part["image_url"]["url"])
                    elif content_part.get("type") == "text":
                        st.markdown(content_part["text"])
                else:
                    st.markdown(content_part)
        else:
            st.markdown(message["content"])

# Chat input
if prompt := st.chat_input("Ask anything... (weather/web tools + optional ICT RAG)"):
    # Prepare user message content
    user_message_content = []

    routing = route_tools(prompt, web_search_provider)
    weather_context = routing["weather_context"]
    search_context = routing["search_context"]

    if routing["weather_error"]:
        st.warning(routing["weather_error"])
    elif weather_context:
        st.info(f"Fetched weather data for: {', '.join(routing['weather_cities'])}")

    if routing["search_error"]:
        st.warning(routing["search_error"])
    elif search_context:
        st.info(f"Loaded web search context from {routing['provider_used']} for: {routing['search_query']}")

    ict_context = ""
    memory_summary = summarize_last_5_turns(current_chat["messages"], model=selected_model)
    if use_ict_rag:
        ict_context, ict_sources, ict_info, rewritten_query = build_ict_rag_context(
            prompt,
            top_k=8,
            domain_filter=rag_domain_filter,
            source_filter=rag_source_filter,
            trust_tier=rag_trust_tier,
            score_threshold=config["retrieval_score_threshold"],
            rewrite_model=selected_model,
        )
        if ict_info.get("error"):
            st.warning(f"ICT RAG skipped: {ict_info['error']}")
        else:
            st.info(
                f"ICT RAG loaded {ict_info['accepted_total']}/{ict_info['retrieved_total']} chunks "
                f"(rejected low-quality: {ict_info['rejected_low_quality']}, threshold: {ict_info['score_threshold']})."
            )
            with st.expander("📚 ICT RAG Sources + Retrieval Filters", expanded=False):
                st.markdown(
                    f"**Rewritten query:** `{rewritten_query}`\n\n"
                    f"**Filters:** domain=`{ict_info['filters']['domain']}`, source=`{ict_info['filters']['source']}`, "
                    f"trust_tier=`{ict_info['filters']['trust_tier']}`, exclude_source=`{ict_info['filters']['exclude_source']}`"
                )
                for src in ict_sources:
                    st.markdown(
                        f"- source_id `{src['source_id']}` | doc_id `{src['doc_id']}` | source `{src['source']}` | "
                        f"domain `{src['domain']}` | trust `{src['trust_tier']}` | chunk `{src['chunk_id']}` | score `{src['score']:.4f}`"
                    )

    # Add text
    user_message_content.append({"type": "text", "text": prompt})
    
    # Add uploaded files
    file_context = ""
    for file_data in st.session_state.uploaded_files:
        if file_data["type"] == "image":
            # Add image to message
            user_message_content.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/png;base64,{file_data['base64']}"
                }
            })
        else:
            # Add document content as
            file_context += f"\n{file_data['name']}: {file_data['content'][:500]}"    
    if file_context:
        user_message_content.append({"type": "text", "text": f"\n\nFiles: {file_context}"})

    if weather_context:
        user_message_content.append({"type": "text", "text": f"\n\nLive weather context:\n{weather_context}"})

    if search_context:
        user_message_content.append({"type": "text", "text": f"\n\nWeb search context:\n{search_context}\n\nUse this context with citations when relevant."})

    if ict_context:
        user_message_content.append({
            "type": "text",
            "text": (
                "\n\nICT Concept RAG context (use this as primary evidence and cite sources):\n"
                f"{ict_context}"
            ),
        })

    if memory_summary:
        user_message_content.append({
            "type": "text",
            "text": f"\n\nConversation summary (last 5 turns only; avoid unrelated topic leakage):\n{memory_summary}",
        })

    current_chat["messages"].append({"role": "user", "content": user_message_content if len(user_message_content) > 1 else prompt})
    
    with st.chat_message("user"):
        st.markdown(prompt)
    
    with st.chat_message("assistant"):
        try:
            llm_messages = [
                {
                    "role": "system",
                    "content": (
                        "You are a production-grade assistant. Prioritize trusted ICT context when provided. "
                        "Do not mix unrelated prior topics. If context is insufficient, state uncertainty."
                    ),
                },
                {"role": "user", "content": user_message_content if len(user_message_content) > 1 else prompt},
            ]
            response = stream_chat_completion(selected_model, llm_messages, temperature)
            full_response = ""
            placeholder = st.empty()
            for chunk in response:
                    if chunk.choices[0].delta.content:                    full_response += chunk.choices[0].delta.content
                    placeholder.markdown(full_response + "▌")
            placeholder.markdown(full_response)
            current_chat["messages"].append({"role": "assistant", "content": full_response})
            pinecone_status = store_conversation_in_pinecone(
                st.session_state.current_chat_id,
                prompt,
                full_response,
                summary=memory_summary,
            )
            if not pinecone_status["success"]:
                st.sidebar.warning(f"Pinecone store skipped: {pinecone_status['reason']}")
                st.sidebar.caption(f"Pinecone status: {pinecone_status}")
        except Exception as e:
            st.error(f"Error: {str(e)}")
    st.rerun()
